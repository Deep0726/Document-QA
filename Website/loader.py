import os
import io
import pandas as pd
from PIL import Image
from google import genai as genai
from langchain_community.document_loaders import UnstructuredFileLoader
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
import fitz
from pptx import Presentation
from docx import Document as DocxDocument

Gemini_API_Key = os.getenv("Gemini_API_Key")
Gemini_API_Key = "AIzaSyAdLBg4GSdrpEsb_Rp8ONKhUng1FK6WBik"
client = genai.Client(api_key=Gemini_API_Key)


# ============================================================
# Splitters — File type ke hisaab se alag alag
# ============================================================

# PDF: Long paragraphs + tables hote hain, thoda bada chunk better hai
# Research: 512-1024 tokens best for mixed factoid + analytical queries
# 1 token ~ 4 chars, so 1000 chars ~ 250 tokens (safe for context window)
pdf_splitter = RecursiveCharacterTextSplitter(
    chunk_size=1000,
    chunk_overlap=200,
    separators=["\n\n", "\n", ".", " ", ""]
)

# DOCX: Heading-wise sections already loader mein bante hain
# Sections moderate size ke hote hain, chhota chunk better retrieval deta hai
docx_splitter = RecursiveCharacterTextSplitter(
    chunk_size=700,
    chunk_overlap=150,
    separators=["\n\n", "\n", ".", " ", ""]
)

# TXT / CSV / MD: Plain text, no inherent structure
# 500 chars ~ 125 tokens, fast retrieval ke liye small chunks
txt_splitter = RecursiveCharacterTextSplitter(
    chunk_size=500,
    chunk_overlap=100,
    separators=["\n\n", "\n", ".", " ", ""]
)

# PPTX: Slide = 1 logical unit, split nahi karenge
# Excel: Row chunks + aggregate chunks, split nahi karenge


# ============================================================
# OCR Model — Gemini
# ============================================================

def ocr_model(img):
    try:
        response = client.models.generate_content(
            model='gemini-2.5-flash',
            contents=[
                "Extract ONLY the text from this image. Do not add explanations. Do not say 'Here is the extracted text'. Preserve formatting.",
                img
            ]
        )
        return response.text if response.text else ""
    except Exception as e:
        print(f"[OCR ERROR]: {e}")
        return ""


# ============================================================
# PDF Reader + Chunking
# ============================================================

def pdf_reader(file_path):
    page_text = []
    doc = fitz.open(file_path)

    for page_num in range(len(doc)):
        page = doc[page_num]

        # Step 1: Normal text extract
        text = page.get_text()

        # Step 2: Dict mode se table structure preserve karo
        blocks = page.get_text("dict")
        for block in blocks.get("blocks", []):
            if block.get("type") == 0:
                for line in block.get("lines", []):
                    line_text = " | ".join(
                        span["text"] for span in line.get("spans", [])
                    )
                    if line_text.strip() and line_text not in text:
                        text += "\n" + line_text

        # Step 3: Images — OCR
        img_list = page.get_images(full=True)
        if len(img_list) >= 1:
            for img_ind, img in enumerate(img_list):
                xref = img[0]
                base_img = doc.extract_image(xref)
                image_byte = base_img["image"]

                if len(image_byte) < 5000:
                    continue

                image = Image.open(io.BytesIO(image_byte))
                img_text = ocr_model(image)

                if img_text:
                    text += "\n [IMAGE]:" + img_text

        page_text.append(Document(
            page_content=text,
            metadata={"page no": page_num}
        ))

    # ✅ PDF-specific splitter se chunk karo
    split_docs = pdf_splitter.split_documents(page_text)
    return split_docs


# ============================================================
# DOCX Reader + Chunking
# ============================================================

def docx_reader(file_path):
    try:
        doc = DocxDocument(file_path)
        docs = []
        current_text = ""
        current_heading = ""

        for para in doc.paragraphs:
            if not para.text.strip():
                continue

            style_name = ""
            if para.style and para.style.name:
                style_name = para.style.name

            if style_name.startswith("Heading"):
                if current_text.strip():
                    docs.append(Document(
                        page_content=f"{current_heading}\n{current_text}".strip(),
                        metadata={"type": "section", "heading": current_heading}
                    ))
                current_heading = para.text.strip()
                current_text = ""
            else:
                current_text += para.text.strip() + "\n"

        if current_text.strip():
            docs.append(Document(
                page_content=f"{current_heading}\n{current_text}".strip(),
                metadata={"type": "section", "heading": current_heading}
            ))

        # Tables alag chunks mein
        for table in doc.tables:
            table_text = ""
            for row in table.rows:
                row_cells = []
                for cell in row.cells:
                    cell_text = cell.text.strip()
                    if cell_text:
                        row_cells.append(cell_text)
                if row_cells:
                    table_text += " | ".join(row_cells) + "\n"
            if table_text.strip():
                docs.append(Document(
                    page_content="[TABLE]:\n" + table_text,
                    metadata={"type": "table"}
                ))

        # Images OCR
        for rel in doc.part.rels.values():
            if "image" in rel.reltype:
                try:
                    image_data = rel.target_part.blob
                    if len(image_data) < 5000:
                        continue
                    image = Image.open(io.BytesIO(image_data))
                    img_text = ocr_model(image)
                    if img_text:
                        docs.append(Document(
                            page_content="[IMAGE]: " + img_text,
                            metadata={"type": "image"}
                        ))
                except Exception as e:
                    print(f"[DOCX IMAGE ERROR]: {e}")
                    continue

        # ✅ DOCX-specific splitter se chunk karo
        split_docs = docx_splitter.split_documents(docs)
        return split_docs

    except Exception as e:
        print(f"[DOCX ERROR]: {e}")
        loader = UnstructuredFileLoader(file_path, mode="elements")
        raw_doc = loader.load()
        docs = []
        current_title = ""
        for doc in raw_doc:
            text = doc.page_content.strip()
            category = doc.metadata.get('category', 'unknown')
            if category in ["Header", "Footer"]:
                continue
            if not text:
                continue
            if "Title" in category:
                current_title = text
                continue
            content = f"{current_title}\n{text}" if current_title else text
            docs.append(Document(
                page_content=content,
                metadata={"type": category}
            ))
        # Fallback mein bhi split karo
        split_docs = docx_splitter.split_documents(docs)
        return split_docs


# ============================================================
# PPT Reader — No splitting (slide = 1 unit)
# ============================================================

def ppt_reader(file_path):
    prs = Presentation(file_path)
    docs = []

    for slide_num, slide in enumerate(prs.slides):
        slide_text = ""
        title = ""

        if slide.shapes.title:
            title = slide.shapes.title.text.strip()

        for shape in slide.shapes:
            if shape == slide.shapes.title:
                continue

            if hasattr(shape, "text"):
                content = shape.text.strip()
                if content:
                    slide_text += "- " + content + "\n"

            if shape.has_table:
                table_text = ""
                for row in shape.table.rows:
                    row_cells = []
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        if cell_text:
                            row_cells.append(cell_text)
                    if row_cells:
                        table_text += " | ".join(row_cells) + "\n"
                if table_text.strip():
                    slide_text += "\n[TABLE]:\n" + table_text

        for idx, shape in enumerate(slide.shapes):
            if shape.shape_type == 13:
                image_bytes = shape.image.blob
                if len(image_bytes) < 5000:
                    continue
                image = Image.open(io.BytesIO(image_bytes))
                img_text = ocr_model(image)
                if img_text:
                    slide_text += f"\n[DIAGRAM {idx+1}]: {img_text}"

        if not slide_text.strip():
            continue

        final_text = f"{title}\n{slide_text}" if title else slide_text

        docs.append(Document(
            page_content=final_text.strip(),
            metadata={"slide": slide_num + 1}
        ))

    # ✅ PPTX: Slide already ek unit hai, split nahi karte
    return docs


# ============================================================
# Excel Reader — Row chunks + Aggregate chunks + DataFrame
# Returns: (docs, dataframes_dict)
# dataframes_dict = {sheet_name: df} — pandas execution ke liye
# ============================================================

def excel_reader(file_path):
    xls = pd.ExcelFile(file_path)
    docs = []
    dataframes = {}  # sheet_name -> cleaned DataFrame

    for sheet_name in xls.sheet_names:
        df = pd.read_excel(file_path, sheet_name=sheet_name)

        # Unnamed columns drop
        df = df.loc[:, ~df.columns.str.contains('^Unnamed')]

        # Empty rows drop
        df = df.dropna(how='all')

        if df.empty:
            continue

        # Column names normalize
        df.columns = [str(col).strip() for col in df.columns]

        # String columns ki values bhi strip karo
        for col in df.select_dtypes(include='object').columns:
            df[col] = df[col].astype(str).str.strip()
            df[col] = df[col].replace('nan', pd.NA)

        # ✅ DataFrame save karo pandas execution ke liye
        dataframes[sheet_name] = df

        # ── Chunk 1: Basic summary ──
        # Columns, row count, unique values — LLM ko structure pata chalega
        summary_parts = [
            f"Sheet '{sheet_name}' has {len(df)} rows and {len(df.columns)} columns.",
            f"Columns: {', '.join(df.columns)}."
        ]
        for col in df.columns:
            unique_vals = df[col].dropna().unique()
            if len(unique_vals) <= 20:
                summary_parts.append(f"Column '{col}' unique values: {list(unique_vals)}.")
            else:
                summary_parts.append(f"Column '{col}' has {len(unique_vals)} unique values.")

        docs.append(Document(
            page_content=" ".join(summary_parts),
            metadata={"sheet": sheet_name, "type": "summary"}
        ))

        # ── Chunk 2: Numeric column stats ──
        # Sum, mean, min, max — kisi bhi numeric data ke liye useful
        num_cols = df.select_dtypes(include='number').columns.tolist()
        if num_cols:
            num_parts = [f"Numeric statistics for sheet '{sheet_name}':"]
            for col in num_cols:
                col_data = df[col].dropna()
                if len(col_data) > 0:
                    num_parts.append(
                        f"'{col}': count={len(col_data)}, "
                        f"sum={round(col_data.sum(), 2)}, "
                        f"mean={round(col_data.mean(), 2)}, "
                        f"min={col_data.min()}, max={col_data.max()}."
                    )
            docs.append(Document(
                page_content=" ".join(num_parts),
                metadata={"sheet": sheet_name, "type": "numeric_stats"}
            ))

        # ── Chunk 3: Categorical value counts ──
        # Har string column ke liye value counts — kisi bhi file pe kaam karega
        cat_cols = df.select_dtypes(include='object').columns.tolist()
        for col in cat_cols:
            col_data = df[col].dropna()
            if len(col_data) == 0:
                continue
            vc = col_data.str.strip().str.lower().value_counts()
            vc_parts = [f"Value counts for column '{col}' in sheet '{sheet_name}':"]
            for val, cnt in vc.items():
                vc_parts.append(f"'{val}' = {cnt} times.")
            docs.append(Document(
                page_content=" ".join(vc_parts),
                metadata={"sheet": sheet_name, "type": "value_counts", "col": col}
            ))

        # ── Chunk 4: Cross-tab chunks ──
        # Group A ke hisaab se Group B ka breakdown
        # Kisi bhi 2 categorical columns pe kaam karega
        if len(cat_cols) >= 2:
            for i in range(len(cat_cols)):
                for j in range(i + 1, len(cat_cols)):
                    col_a = cat_cols[i]
                    col_b = cat_cols[j]
                    try:
                        temp = df[[col_a, col_b]].dropna()
                        if len(temp) == 0:
                            continue
                        temp = temp.copy()
                        temp[col_a] = temp[col_a].str.strip().str.lower()
                        temp[col_b] = temp[col_b].str.strip().str.lower()
                        cross = temp.groupby([col_a, col_b]).size().reset_index(name='count')
                        cross_parts = [
                            f"Breakdown of '{col_b}' grouped by '{col_a}' in sheet '{sheet_name}':"
                        ]
                        for _, r in cross.iterrows():
                            cross_parts.append(
                                f"{col_a}='{r[col_a]}' with {col_b}='{r[col_b]}': {r['count']} records."
                            )
                        docs.append(Document(
                            page_content=" ".join(cross_parts),
                            metadata={
                                "sheet": sheet_name,
                                "type": "crosstab",
                                "cols": f"{col_a}x{col_b}"
                            }
                        ))
                    except Exception as e:
                        print(f"[EXCEL CROSSTAB ERROR] {col_a}x{col_b}: {e}")
                        continue

        # ── Chunk 5: Individual row chunks ──
        # Specific row lookup ke liye (RAG)
        for i, row in df.iterrows():
            row_text = ""
            for col in df.columns:
                value = row[col]
                if pd.notna(value) and str(value).strip() not in ('', 'nan', 'None'):
                    row_text += f"{col}: {str(value).strip()}. "
            if row_text.strip():
                docs.append(Document(
                    page_content=row_text.strip(),
                    metadata={"sheet": sheet_name, "row": i}
                ))

    return docs, dataframes


# ============================================================
# Image Reader — No splitting (single OCR output)
# ============================================================

def image_reader(file_path):
    image = Image.open(file_path)
    text = ocr_model(image)

    docs = [
        Document(
            page_content=text.strip() if text else "No text extracted",
            metadata={"type": "image"}
        )
    ]
    return docs


# ============================================================
# TXT / CSV / MD Reader + Chunking
# ============================================================

def txt_reader(file_path):
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        text = f.read()

    docs = [Document(page_content=text, metadata={"source": file_path})]
    split_docs = txt_splitter.split_documents(docs)
    return split_docs


# ============================================================
# Universal Loader
# Returns:
#   - Excel: (docs, dataframes_dict)
#   - All others: docs list
# ============================================================

def load_document(file_path):
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        return pdf_reader(file_path)

    elif ext in [".docx", ".doc"]:
        return docx_reader(file_path)

    elif ext in [".pptx", ".ppt"]:
        return ppt_reader(file_path)

    elif ext in [".xlsx", ".xls"]:
        return excel_reader(file_path)  # Returns (docs, dataframes)

    elif ext in [".png", ".jpg", ".jpeg", ".webp", ".bmp"]:
        return image_reader(file_path)

    elif ext in [".txt", ".md", ".csv"]:
        return txt_reader(file_path)

    else:
        raise ValueError(f"Unsupported file type: {ext}")